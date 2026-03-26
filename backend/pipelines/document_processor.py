"""
Document Processor - Multi-format document to page images conversion.

Supports:
- Office: DOCX, PPTX, XLSX
- Text: TXT, MD, CSV, JSON, HTML
- Images: PNG, JPG, JPEG, WEBP, TIFF, BMP, GIF
- PDF: via existing PDFProcessor

Common interface: DocumentProcessor.to_pages(file_path) -> List[PIL.Image]
"""
import os
import io
import json
import csv
import logging
import textwrap
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# Page rendering constants
PAGE_WIDTH = 1700
PAGE_HEIGHT = 2200
MARGIN = 80
FONT_SIZE = 28
LINE_HEIGHT = 36
TITLE_FONT_SIZE = 40
HEADER_HEIGHT = 120
BG_COLOR = (255, 255, 255)
TEXT_COLOR = (30, 30, 30)
LIGHT_GRAY = (200, 200, 200)
HEADER_BG = (245, 245, 245)


def _get_font(size: int = FONT_SIZE, bold: bool = False) -> ImageFont.FreeTypeFont:
    """Get a font, falling back to default if system fonts not available."""
    font_paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ]
    for fp in font_paths:
        if fp and os.path.exists(fp):
            try:
                return ImageFont.truetype(fp, size)
            except Exception:
                continue
    return ImageFont.load_default()


def _render_text_to_pages(
    text: str,
    title: str = "",
    max_chars_per_line: int = 80
) -> List[Image.Image]:
    """
    Render plain text content into page images.
    Splits text into pages that fit within page dimensions.
    """
    font = _get_font(FONT_SIZE)
    title_font = _get_font(TITLE_FONT_SIZE, bold=True)

    # Wrap lines
    lines = []
    for raw_line in text.split("\n"):
        if raw_line.strip() == "":
            lines.append("")
        else:
            wrapped = textwrap.wrap(raw_line, width=max_chars_per_line)
            lines.extend(wrapped if wrapped else [""])

    # Calculate lines per page
    usable_height = PAGE_HEIGHT - 2 * MARGIN - (HEADER_HEIGHT if title else 0)
    lines_per_page = max(1, usable_height // LINE_HEIGHT)

    # Split into page chunks
    pages = []
    for i in range(0, max(1, len(lines)), lines_per_page):
        page_lines = lines[i:i + lines_per_page]

        img = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)
        draw = ImageDraw.Draw(img)

        y = MARGIN

        # Title on first page only
        if i == 0 and title:
            draw.rectangle(
                [(0, 0), (PAGE_WIDTH, HEADER_HEIGHT)],
                fill=HEADER_BG
            )
            draw.text((MARGIN, 30), title, fill=TEXT_COLOR, font=title_font)
            y = HEADER_HEIGHT + 20

        # Draw text lines
        for line in page_lines:
            draw.text((MARGIN, y), line, fill=TEXT_COLOR, font=font)
            y += LINE_HEIGHT

        # Page number
        page_num_text = f"Page {len(pages) + 1}"
        draw.text(
            (PAGE_WIDTH - MARGIN - 100, PAGE_HEIGHT - 50),
            page_num_text,
            fill=LIGHT_GRAY,
            font=_get_font(20)
        )

        pages.append(img)

    return pages if pages else [Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)]


def _render_table_to_pages(
    headers: List[str],
    rows: List[List[str]],
    title: str = ""
) -> List[Image.Image]:
    """
    Render tabular data (CSV/XLSX) into page images with a table layout.
    """
    font = _get_font(FONT_SIZE - 4)
    title_font = _get_font(TITLE_FONT_SIZE, bold=True)
    header_font = _get_font(FONT_SIZE - 2, bold=True)

    # Calculate column widths
    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    max_col_width = (PAGE_WIDTH - 2 * MARGIN) // max(1, num_cols)
    col_width = min(max_col_width, 300)
    table_width = col_width * num_cols
    row_height = 40

    # Rows per page
    usable_height = PAGE_HEIGHT - 2 * MARGIN - HEADER_HEIGHT - row_height  # header row
    rows_per_page = max(1, usable_height // row_height)

    pages = []
    for chunk_start in range(0, max(1, len(rows)), rows_per_page):
        chunk_rows = rows[chunk_start:chunk_start + rows_per_page]

        img = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)
        draw = ImageDraw.Draw(img)

        y = MARGIN

        # Title
        if chunk_start == 0 and title:
            draw.rectangle([(0, 0), (PAGE_WIDTH, HEADER_HEIGHT)], fill=HEADER_BG)
            draw.text((MARGIN, 30), title, fill=TEXT_COLOR, font=title_font)
            y = HEADER_HEIGHT + 20

        # Column headers
        if headers:
            x = MARGIN
            draw.rectangle(
                [(MARGIN, y), (MARGIN + table_width, y + row_height)],
                fill=(230, 230, 240)
            )
            for j, h in enumerate(headers[:num_cols]):
                cell_text = str(h)[:col_width // 8]
                draw.text((x + 5, y + 8), cell_text, fill=TEXT_COLOR, font=header_font)
                x += col_width
            y += row_height

        # Data rows
        for row_idx, row in enumerate(chunk_rows):
            x = MARGIN
            bg = BG_COLOR if row_idx % 2 == 0 else (248, 248, 248)
            draw.rectangle(
                [(MARGIN, y), (MARGIN + table_width, y + row_height)],
                fill=bg
            )
            # Grid lines
            draw.line([(MARGIN, y), (MARGIN + table_width, y)], fill=LIGHT_GRAY)

            for j in range(num_cols):
                cell_val = str(row[j]) if j < len(row) else ""
                cell_text = cell_val[:col_width // 8]
                draw.text((x + 5, y + 8), cell_text, fill=TEXT_COLOR, font=font)
                # Vertical line
                draw.line([(x, y), (x, y + row_height)], fill=LIGHT_GRAY)
                x += col_width
            # Right edge
            draw.line([(x, y), (x, y + row_height)], fill=LIGHT_GRAY)
            y += row_height

        # Bottom border
        draw.line([(MARGIN, y), (MARGIN + table_width, y)], fill=LIGHT_GRAY)

        # Page number
        page_num_text = f"Page {len(pages) + 1}"
        draw.text(
            (PAGE_WIDTH - MARGIN - 100, PAGE_HEIGHT - 50),
            page_num_text,
            fill=LIGHT_GRAY,
            font=_get_font(20)
        )

        pages.append(img)

    return pages if pages else [Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)]


class DocxProcessor:
    """Convert DOCX files to page images."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        from docx import Document

        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            text_parts.append(para.text)

        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                text_parts.append(row_text)

        full_text = "\n".join(text_parts)
        title = Path(file_path).stem

        logger.info(f"DOCX: extracted {len(text_parts)} paragraphs/rows from {title}")
        return _render_text_to_pages(full_text, title=title)


class PptxProcessor:
    """Convert PPTX files to slide images."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu

        prs = Presentation(file_path)
        slide_width = prs.slide_width or Emu(9144000)  # default 10"
        slide_height = prs.slide_height or Emu(6858000)  # default 7.5"

        pages = []
        for slide_idx, slide in enumerate(prs.slides):
            # Extract all text from slide
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        texts.append(row_text)

            # Render slide as image
            img = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)
            draw = ImageDraw.Draw(img)

            # Slide header
            draw.rectangle([(0, 0), (PAGE_WIDTH, HEADER_HEIGHT)], fill=(50, 50, 80))
            title_font = _get_font(TITLE_FONT_SIZE, bold=True)
            slide_title = texts[0] if texts else f"Slide {slide_idx + 1}"
            # Truncate long titles
            if len(slide_title) > 60:
                slide_title = slide_title[:57] + "..."
            draw.text(
                (MARGIN, 35),
                slide_title,
                fill=(255, 255, 255),
                font=title_font
            )

            # Body text
            font = _get_font(FONT_SIZE)
            y = HEADER_HEIGHT + 40
            body_texts = texts[1:] if len(texts) > 1 else texts
            for line in body_texts:
                wrapped = textwrap.wrap(line, width=80)
                for wl in wrapped:
                    if y > PAGE_HEIGHT - MARGIN - 60:
                        break
                    draw.text((MARGIN, y), wl, fill=TEXT_COLOR, font=font)
                    y += LINE_HEIGHT

            # Slide number
            draw.text(
                (PAGE_WIDTH - MARGIN - 100, PAGE_HEIGHT - 50),
                f"Slide {slide_idx + 1}",
                fill=LIGHT_GRAY,
                font=_get_font(20)
            )

            pages.append(img)

        logger.info(f"PPTX: rendered {len(pages)} slides from {Path(file_path).name}")
        return pages if pages else [Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)]


class XlsxProcessor:
    """Convert XLSX files to table images."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        all_pages = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            headers = None

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                str_row = [str(cell) if cell is not None else "" for cell in row]
                if row_idx == 0:
                    headers = str_row
                else:
                    rows.append(str_row)

            if not headers and not rows:
                continue

            title = f"{Path(file_path).stem} — {sheet_name}"
            pages = _render_table_to_pages(
                headers=headers or [],
                rows=rows,
                title=title
            )
            all_pages.extend(pages)

        wb.close()
        logger.info(f"XLSX: rendered {len(all_pages)} pages from {Path(file_path).name}")
        return all_pages if all_pages else [Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)]


class TextFileProcessor:
    """Convert text-based files (TXT, MD, HTML) to page images."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        ext = Path(file_path).suffix.lower()
        title = Path(file_path).stem

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        if ext == ".html":
            # Strip HTML tags for text rendering
            import re
            content = re.sub(r'<[^>]+>', '', content)
            content = content.strip()

        logger.info(f"TEXT: {len(content)} chars from {Path(file_path).name}")
        return _render_text_to_pages(content, title=title)


class CsvProcessor:
    """Convert CSV/TSV files to table images."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        ext = Path(file_path).suffix.lower()
        delimiter = "\t" if ext == ".tsv" else ","

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            all_rows = list(reader)

        if not all_rows:
            return [Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BG_COLOR)]

        headers = all_rows[0]
        rows = all_rows[1:]
        title = Path(file_path).stem

        logger.info(f"CSV: {len(rows)} rows, {len(headers)} cols from {Path(file_path).name}")
        return _render_table_to_pages(headers=headers, rows=rows, title=title)


class JsonProcessor:
    """Convert JSON files to page images (text or table depending on structure)."""

    @staticmethod
    def to_pages(file_path: str) -> List[Image.Image]:
        title = Path(file_path).stem

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)

        # If JSON is a list of dicts, render as table
        if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
            headers = list(data[0].keys())
            rows = []
            for item in data:
                row = [str(item.get(h, "")) for h in headers]
                rows.append(row)
            logger.info(f"JSON (table): {len(rows)} rows from {Path(file_path).name}")
            return _render_table_to_pages(headers=headers, rows=rows, title=title)
        else:
            # Render as formatted text
            formatted = json.dumps(data, indent=2, ensure_ascii=False)
            logger.info(f"JSON (text): {len(formatted)} chars from {Path(file_path).name}")
            return _render_text_to_pages(formatted, title=title)


class DocumentProcessor:
    """
    Universal document processor.
    Common interface: DocumentProcessor.to_pages(file_path) -> List[PIL.Image]
    """

    # Extension to processor mapping
    PROCESSORS = {
        # Office
        ".docx": DocxProcessor,
        ".doc": DocxProcessor,  # Best effort via python-docx
        ".pptx": PptxProcessor,
        ".xlsx": XlsxProcessor,
        ".xls": XlsxProcessor,
        # Text
        ".txt": TextFileProcessor,
        ".md": TextFileProcessor,
        ".html": TextFileProcessor,
        ".htm": TextFileProcessor,
        ".log": TextFileProcessor,
        ".xml": TextFileProcessor,
        ".yaml": TextFileProcessor,
        ".yml": TextFileProcessor,
        ".toml": TextFileProcessor,
        ".ini": TextFileProcessor,
        ".cfg": TextFileProcessor,
        ".py": TextFileProcessor,
        ".js": TextFileProcessor,
        ".ts": TextFileProcessor,
        ".java": TextFileProcessor,
        ".go": TextFileProcessor,
        ".rs": TextFileProcessor,
        ".c": TextFileProcessor,
        ".cpp": TextFileProcessor,
        ".h": TextFileProcessor,
        # Data
        ".csv": CsvProcessor,
        ".tsv": CsvProcessor,
        ".json": JsonProcessor,
    }

    # Image extensions handled by visual_engine directly
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp", ".gif"}

    @classmethod
    def supported_extensions(cls) -> set:
        """Get all supported file extensions."""
        return set(cls.PROCESSORS.keys()) | cls.IMAGE_EXTENSIONS | {".pdf"}

    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if a file type is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in cls.supported_extensions()

    @classmethod
    def to_pages(cls, file_path: str) -> List[Image.Image]:
        """
        Convert any supported document to a list of page images.

        Args:
            file_path: Path to the document file

        Returns:
            List of PIL Image objects (one per page/slide/sheet)

        Raises:
            ValueError: If file type is not supported
        """
        ext = Path(file_path).suffix.lower()

        # PDF — delegate to existing PDFProcessor
        if ext == ".pdf":
            from backend.pipelines.visual_engine import PDFProcessor
            return PDFProcessor.convert_to_images(file_path)

        # Images — direct load
        if ext in cls.IMAGE_EXTENSIONS:
            return [Image.open(file_path).convert("RGB")]

        # Document processors
        processor = cls.PROCESSORS.get(ext)
        if processor:
            return processor.to_pages(file_path)

        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {', '.join(sorted(cls.supported_extensions()))}"
        )

    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """
        Extract raw text from a document (for Text Agent indexing).
        Returns text content without rendering to images.
        """
        ext = Path(file_path).suffix.lower()

        if ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
            return "\n".join(parts)

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                parts.append(para.text.strip())
            return "\n".join(parts)

        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            return "\n".join(parts)

        elif ext in {".csv", ".tsv"}:
            delimiter = "\t" if ext == ".tsv" else ","
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f, delimiter=delimiter)
                return "\n".join(" ".join(row) for row in reader)

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)

        elif ext == ".html" or ext == ".htm":
            import re
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return re.sub(r'<[^>]+>', '', content).strip()

        elif ext in cls.PROCESSORS:
            # Generic text file
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        return ""
